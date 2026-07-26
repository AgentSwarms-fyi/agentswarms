"""
Native PowerPoint renderer (python-pptx) for AgentSwarms server-side doc-gen.

Takes a filled plan (the same PptxPlan the browser planner produces, with chart
categories/series already computed and each diagram slide carrying a pre-rendered
SVG) and produces a .pptx with NATIVE, EDITABLE charts, tables, KPI cards and
text — plus diagrams embedded as images (rasterised from the client's SVG).

Design mirrors the browser builder (deep-ink cover, section dividers, KPI cards,
chart cards, key-insight bars) but follows the Anthropic pptx skill's guidance:
no "accent line under the title" (an AI hallmark), generous margins, size
contrast, every slide carries a visual.
"""
from __future__ import annotations

import io
import re
from typing import Any

import cairosvg
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

EMU_PER_IN = 914400
CW = 13.333
CH = 7.5
M = 0.6
CONTENT_W = CW - M * 2

INK = "0F172A"
INK_DEEP = "0B1220"
SUB = "64748B"
BODY = "334155"
BORDER = "E7EBF0"
CARD = "F8FAFC"
DEFAULT_ACCENT = "4F46E5"
CHART_PALETTE = ["4F46E5", "0EA5E9", "10B981", "F59E0B", "EF4444", "8B5CF6", "EC4899", "14B8A6"]


def _hex(c: str | None, fallback: str = DEFAULT_ACCENT) -> str:
    h = re.sub(r"[^0-9a-fA-F]", "", (c or ""))
    return h.upper() if len(h) == 6 else fallback


def _rgb(c: str) -> RGBColor:
    return RGBColor.from_string(_hex(c))


def _lum(c: str) -> float:
    n = int(_hex(c), 16)
    return (0.2126 * ((n >> 16) & 255) + 0.7152 * ((n >> 8) & 255) + 0.0722 * (n & 255)) / 255


def _on_light(c: str) -> str:
    h = _hex(c)
    for _ in range(6):
        if _lum(h) <= 0.6:
            break
        n = int(h, 16)
        h = "".join(f"{round(((n >> s) & 255) * 0.72):02X}" for s in (16, 8, 0))
    return h


def _compact(n: float) -> str:
    a = abs(n)
    if a >= 1e9:
        return f"{n/1e9:.1f}".rstrip("0").rstrip(".") + "B"
    if a >= 1e6:
        return f"{n/1e6:.1f}".rstrip("0").rstrip(".") + "M"
    if a >= 1e3:
        return f"{n/1e3:.1f}".rstrip("0").rstrip(".") + "K"
    return str(int(n)) if float(n).is_integer() else f"{n:.2f}".rstrip("0").rstrip(".")


class Deck:
    def __init__(self, plan: dict[str, Any]):
        self.plan = plan
        self.accent = _hex(plan.get("accent"), DEFAULT_ACCENT)
        self.accent_ink = _on_light(self.accent)
        self.prs = Presentation()
        self.prs.slide_width = Inches(CW)
        self.prs.slide_height = Inches(CH)
        self.blank = self.prs.slide_layouts[6]

    # ── low-level helpers ────────────────────────────────────────────────
    def _slide(self, bg: str | None = None):
        s = self.prs.slides.add_slide(self.blank)
        if bg:
            s.background.fill.solid()
            s.background.fill.fore_color.rgb = _rgb(bg)
        return s

    def _rect(self, slide, x, y, w, h, fill, *, radius=False, line=None, line_w=1.0):
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
        if line:
            shp.line.color.rgb = _rgb(line)
            shp.line.width = Pt(line_w)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, x, y, w, h, runs, *, size=14, color=BODY, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None, shrink=False):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        items = runs if isinstance(runs, list) else [(runs, {})]
        first = True
        for text, opt in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = opt.get("align", align)
            if opt.get("space_after") is not None:
                p.space_after = Pt(opt["space_after"])
            if opt.get("bullet"):
                _set_bullet(p)
            r = p.add_run()
            r.text = text
            r.font.size = Pt(opt.get("size", size))
            r.font.bold = opt.get("bold", bold)
            r.font.name = "Segoe UI"
            r.font.color.rgb = _rgb(opt.get("color", color))
            if spacing:
                _set_char_spacing(r, spacing)
        return tb

    def _svg_image(self, slide, svg: str, x, y, w, h):
        png = cairosvg.svg2png(
            bytestring=svg.encode("utf-8"),
            output_width=int(w * 150),
            output_height=int(h * 150),
        )
        slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(y), Inches(w), Inches(h))

    # ── deck ─────────────────────────────────────────────────────────────
    def build(self) -> bytes:
        self._cover()
        section_no = 0
        current_section = ""
        for s in self.plan.get("slides", []):
            layout = self._layout(s)
            if layout == "section":
                section_no += 1
                current_section = s.get("title") or f"Section {section_no}"
                self._section(section_no, s)
                continue
            slide = self._slide("FFFFFF")
            # accent spine
            self._rect(slide, 0, 0, 0.16, CH, self.accent)
            self._header(slide, s.get("title", ""), current_section or self.plan.get("title", ""))
            bottom = 6.25 if s.get("takeaway") else 6.95
            top = 1.55
            self._content(slide, s, layout, top, bottom)
            if s.get("takeaway"):
                self._takeaway(slide, s["takeaway"])
            if s.get("notes"):
                slide.notes_slide.notes_text_frame.text = s["notes"]
        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()

    def _layout(self, s: dict) -> str:
        if s.get("layout"):
            return s["layout"]
        if s.get("diagramSvg") or s.get("diagram"):
            return "diagram"
        if s.get("kpis"):
            return "kpi"
        if _chart_has_data(s.get("chart")):
            return "chart"
        if s.get("table"):
            return "table"
        return "bullets"

    def _cover(self):
        s = self._slide(INK_DEEP)
        self._rect(s, 0.9, 2.02, 0.62, 0.09, self.accent)
        self._text(s, 1.64, 1.82, 9, 0.4, "REPORT", size=11.5, bold=True,
                   color=_light(self.accent), spacing=3)
        self._text(s, 0.85, 2.45, 10.5, 2.3, self.plan.get("title", "Untitled"),
                   size=44, bold=True, color="FFFFFF")
        if self.plan.get("subtitle"):
            self._text(s, 0.9, 4.95, 9.5, 1.0, self.plan["subtitle"], size=18, color="CBD5E1")
        self._rect(s, 0.9, 6.45, 3.0, 0.02, "334155")

    def _section(self, no: int, s: dict):
        slide = self._slide(INK_DEEP)
        self._text(slide, 6.4, 1.1, 6.6, 6.2, f"{no:02d}", size=280, bold=True,
                   color=_light(self.accent), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM)
        self._rect(slide, 0.9, 2.72, 0.62, 0.09, self.accent)
        self._text(slide, 1.64, 2.52, 8, 0.4, f"SECTION {no:02d}", size=11.5, bold=True,
                   color=_light(self.accent), spacing=3)
        self._text(slide, 0.9, 3.02, 9.4, 1.5, s.get("title", ""), size=36, bold=True, color="FFFFFF")
        if s.get("subtitle"):
            self._text(slide, 0.92, 4.6, 8.8, 0.8, s["subtitle"], size=15, color="CBD5E1")
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    def _header(self, slide, title: str, kicker: str):
        if kicker:
            self._text(slide, M, 0.4, CONTENT_W, 0.28, kicker.upper()[:60], size=10.5,
                       bold=True, color=self.accent_ink, spacing=2)
        self._text(slide, M, 0.66, CONTENT_W, 0.7, title, size=24, bold=True, color=INK)
        # NOTE: no accent underline (Anthropic skill flags it as an AI hallmark);
        # a full-width hairline only.
        self._rect(slide, M, 1.32, CONTENT_W, 0.012, BORDER)

    def _takeaway(self, slide, text: str):
        self._rect(slide, M, 6.4, CONTENT_W, 0.58, _tint(self.accent, 0.9), radius=True)
        self._rect(slide, M, 6.4, 0.1, 0.58, self.accent, radius=True)
        self._text(slide, M + 0.32, 6.4, CONTENT_W - 0.6, 0.58,
                   [("KEY INSIGHT   ", {"bold": True, "color": self.accent_ink}),
                    (text, {"color": INK})],
                   size=12, anchor=MSO_ANCHOR.MIDDLE)

    # ── content dispatch ─────────────────────────────────────────────────
    def _content(self, slide, s, layout, top, bottom):
        if layout == "diagram" and s.get("diagramSvg"):
            self._svg_image(slide, s["diagramSvg"], M, top, CONTENT_W, bottom - top)
            return
        if layout == "kpi" and s.get("kpis"):
            self._kpis(slide, s["kpis"][:5], s.get("bullets"), bottom)
            return
        has_chart = _chart_has_data(s.get("chart"))
        has_visual = has_chart or bool(s.get("table"))
        has_text = bool(s.get("bullets") or s.get("paragraph"))
        if has_visual and has_text:
            if has_chart:
                self._chart(slide, s["chart"], M, top, 7.3, bottom - top)
            else:
                self._table(slide, s["table"], M, top, 7.3)
            ty = top + 0.05
            if s.get("paragraph"):
                self._text(slide, 8.2, ty, 4.5, 1.6, s["paragraph"], size=14, color=BODY)
                ty += 1.7
            if s.get("bullets"):
                self._bullets(slide, s["bullets"], 8.2, ty, 4.5, bottom - ty)
        elif has_visual:
            if has_chart:
                self._chart(slide, s["chart"], M, top, CONTENT_W, bottom - top)
            else:
                self._table(slide, s["table"], M, top, CONTENT_W)
        else:
            y = top
            if s.get("paragraph"):
                self._text(slide, M, y, CONTENT_W, 1.6, s["paragraph"], size=16, color=BODY)
                y += 1.7
            if s.get("bullets"):
                self._bullets(slide, s["bullets"], M, y, CONTENT_W, bottom - y)

    def _bullets(self, slide, bullets, x, y, w, h):
        runs = [(b, {"bullet": True, "size": 18, "color": BODY, "space_after": 12}) for b in bullets]
        self._text(slide, x, y, w, h, runs)

    def _card(self, slide, x, y, w, h, top_rule=False):
        self._rect(slide, x, y, w, h, "FFFFFF", radius=True, line=BORDER)
        if top_rule:
            self._rect(slide, x, y, w, 0.09, self.accent)

    def _kpis(self, slide, kpis, bullets, bottom):
        gap = 0.28
        cw = (CONTENT_W - gap * (len(kpis) - 1)) / len(kpis)
        cy, ch = 1.95, 2.3
        for i, k in enumerate(kpis):
            x = M + i * (cw + gap)
            self._card(slide, x, cy, cw, ch, top_rule=True)
            self._text(slide, x + 0.26, cy + 0.34, cw - 0.5, 0.7, str(k.get("label", "")).upper(),
                       size=10.5, bold=True, color=SUB, spacing=1)
            self._text(slide, x + 0.24, cy + 0.9, cw - 0.42, 0.9, str(k.get("value", "")),
                       size=32, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
            if k.get("delta"):
                good = k.get("positive") is not False
                self._rect(slide, x + 0.26, cy + ch - 0.62, 1.2, 0.34,
                           "DCFCE7" if good else "FEE2E2", radius=True)
                self._text(slide, x + 0.26, cy + ch - 0.62, 1.2, 0.34, str(k["delta"]),
                           size=10.5, bold=True, color="047857" if good else "B91C1C",
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if bullets:
            self._bullets(slide, bullets, M, 4.55, CONTENT_W, bottom - 4.55)

    def _table(self, slide, t, x, y, w):
        cols = t.get("columns", [])
        rows = t.get("rows", [])
        if not cols:
            return
        n = len(rows) + 1
        gshape = slide.shapes.add_table(n, len(cols), Inches(x), Inches(y), Inches(w),
                                        Inches(min(0.4 * n, 5.0)))
        table = gshape.table
        for c, col in enumerate(cols):
            cell = table.cell(0, c)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(self.accent)
            for pr in cell.text_frame.paragraphs:
                for rn in pr.runs:
                    rn.font.bold = True
                    rn.font.color.rgb = _rgb("FFFFFF")
                    rn.font.size = Pt(12)
                    rn.font.name = "Segoe UI"
        for r, row in enumerate(rows):
            for c in range(len(cols)):
                val = row[c] if c < len(row) else ""
                cell = table.cell(r + 1, c)
                cell.text = "" if val is None else str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FFFFFF" if r % 2 == 0 else CARD)
                for pr in cell.text_frame.paragraphs:
                    for rn in pr.runs:
                        rn.font.size = Pt(11)
                        rn.font.name = "Segoe UI"
                        rn.font.color.rgb = _rgb(BODY)

    def _chart(self, slide, chart, x, y, w, h):
        self._card(slide, x, y, w, h)
        cd = CategoryChartData()
        cd.categories = chart.get("categories", [])
        series = chart.get("series", []) or []
        for ser in series:
            cd.add_series(ser.get("name", "Series"), tuple(float(v or 0) for v in ser.get("values", [])))
        ctype = {
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "area": XL_CHART_TYPE.AREA,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        }.get(chart.get("type"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = slide.shapes.add_chart(ctype, Inches(x + 0.25), Inches(y + 0.25),
                                    Inches(w - 0.5), Inches(h - 0.5), cd)
        ch = gf.chart
        ch.has_title = False
        is_pie = chart.get("type") in ("pie", "doughnut")
        ch.has_legend = is_pie or len(series) > 1
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
        try:
            plot = ch.plots[0]
            plot.vary_by_categories = is_pie or len(series) == 1
            for idx, ser in enumerate(plot.series):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = _rgb(CHART_PALETTE[idx % len(CHART_PALETTE)])
        except Exception:
            pass


def _light(c: str) -> str:
    """Lighten toward white for readable text on the dark cover/section bg."""
    h = _hex(c)
    for _ in range(6):
        if _lum(h) >= 0.62:
            break
        n = int(h, 16)
        h = "".join(f"{round(v + (255 - v) * 0.35):02X}" for v in ((n >> 16) & 255, (n >> 8) & 255, n & 255))
    return h


def _tint(c: str, amt: float) -> str:
    n = int(_hex(c), 16)
    return "".join(
        f"{round(v + (255 - v) * amt):02X}" for v in ((n >> 16) & 255, (n >> 8) & 255, n & 255)
    )


def _chart_has_data(chart) -> bool:
    if not chart:
        return False
    cats = chart.get("categories") or []
    series = chart.get("series") or []
    return len(cats) > 0 and any(len(s.get("values") or []) > 0 for s in series)


def _set_bullet(paragraph):
    from pptx.oxml.ns import qn
    pPr = paragraph._pPr if paragraph._pPr is not None else paragraph.get_or_add_pPr()
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪"})
    pPr.append(buChar)


def _set_char_spacing(run, pts: float):
    from pptx.oxml.ns import qn
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def render_pptx(plan: dict[str, Any]) -> bytes:
    return Deck(plan).build()
